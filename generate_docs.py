"""Generate Word and PowerPoint versions of the Factory.AI dashboard guide."""

from docx import Document
from docx.shared import Inches, Pt, RGBColor
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.enum.table import WD_TABLE_ALIGNMENT
from pptx import Presentation
from pptx.util import Inches as PptInches, Pt as PptPt
from pptx.dml.color import RGBColor as PptRGB
from pptx.enum.text import PP_ALIGN
import os

OUT_DIR = os.path.dirname(os.path.abspath(__file__))

# ── EY brand colors ──
EY_YELLOW = RGBColor(0xFF, 0xE6, 0x00)
EY_DARK = RGBColor(0x2E, 0x2E, 0x38)
EY_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
EY_GREY = RGBColor(0x74, 0x74, 0x7A)

PPT_YELLOW = PptRGB(0xFF, 0xE6, 0x00)
PPT_DARK = PptRGB(0x2E, 0x2E, 0x38)
PPT_WHITE = PptRGB(0xFF, 0xFF, 0xFF)


# ═══════════════════════════════════════════════════════════════════════════
# WORD DOCUMENT
# ═══════════════════════════════════════════════════════════════════════════

def make_word():
    doc = Document()

    style = doc.styles["Normal"]
    style.font.name = "Calibri"
    style.font.size = Pt(11)
    style.font.color.rgb = EY_DARK

    for level in range(1, 4):
        hs = doc.styles[f"Heading {level}"]
        hs.font.color.rgb = EY_DARK
        hs.font.name = "Calibri"

    # ── Title page ──
    for _ in range(6):
        doc.add_paragraph()

    title = doc.add_paragraph()
    title.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run = title.add_run("Building and Deploying a Live Dashboard\nwith Factory.AI")
    run.font.size = Pt(28)
    run.font.color.rgb = EY_DARK
    run.bold = True

    subtitle = doc.add_paragraph()
    subtitle.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run = subtitle.add_run("A Step-by-Step Guide for EY Professionals")
    run.font.size = Pt(16)
    run.font.color.rgb = EY_GREY

    doc.add_page_break()

    # ── What You Will Build ──
    doc.add_heading("What You Will Build", level=1)
    doc.add_paragraph(
        "A live, interactive dashboard that:"
    )
    items = [
        "Visualizes data from an Excel file you already maintain",
        "Updates automatically when your Excel file changes",
        "Is accessible to everyone at EY via a simple web link",
        "Gives you admin control over who can edit data",
        "Is secured with Microsoft SSO \u2014 users sign in with their real EY account",
    ]
    for item in items:
        doc.add_paragraph(item, style="List Bullet")

    p = doc.add_paragraph()
    run = p.add_run("Time required: ")
    run.bold = True
    p.add_run("~30 minutes (including SSO setup)")

    p = doc.add_paragraph()
    run = p.add_run("Technical experience required: ")
    run.bold = True
    p.add_run("None. Factory.AI builds the code for you.")

    # ── Prerequisites ──
    doc.add_heading("What You Need Before Starting", level=1)

    table = doc.add_table(rows=6, cols=2)
    table.style = "Light Grid Accent 1"
    table.alignment = WD_TABLE_ALIGNMENT.CENTER
    headers = ["Item", "Where to Get It"]
    rows_data = [
        ["Factory.AI CLI (Droid)", "Pre-installed, or visit factory.ai"],
        ["A GitHub account", "github.com (free)"],
        ["A Streamlit Community Cloud account", "streamlit.io/cloud (free, sign in with GitHub)"],
        ["Your data in an Excel file (.xlsx)", "Any Excel file on your machine or OneDrive"],
        ["Azure AD app registration (for SSO)", "Request from your IT team \u2014 see Step 8 and email template below"],
    ]
    for i, h in enumerate(headers):
        cell = table.rows[0].cells[i]
        cell.text = h
        cell.paragraphs[0].runs[0].bold = True
    for r, row_data in enumerate(rows_data):
        for c, val in enumerate(row_data):
            table.rows[r + 1].cells[c].text = val

    # ── Send IT request early ──
    doc.add_paragraph()
    p = doc.add_paragraph()
    run = p.add_run("IMPORTANT: Send the Azure request now. ")
    run.bold = True
    run.font.color.rgb = RGBColor(0xE0, 0x00, 0x4D)
    p.add_run(
        "Step 8 requires your IT team to register an app in Azure, which may take a day or two. "
        "Send them the request now so it is ready by the time you reach Step 8. "
        "The dashboard works without SSO while you wait, but SSO is required before sharing sensitive data."
    )

    doc.add_heading("Email Template \u2014 Send This to Your IT Team Now", level=2)
    email_text = (
        'Hi,\n\n'
        'I am building an internal Streamlit dashboard and need an Azure AD '
        '(Microsoft Entra ID) app registration to enable Microsoft SSO. '
        'Could you help with the following?\n\n'
        '1. Register a new app called "Factory AI Dashboard"\n'
        '2. Set supported account types to "Accounts in this organizational directory only"\n'
        '3. Add a Web redirect URI (I will provide the URL once deployed \u2014 '
        'placeholder: https://placeholder.streamlit.app)\n'
        '4. Create a client secret\n'
        '5. Add the Microsoft Graph "User.Read" delegated permission and grant admin consent\n'
        '6. Send me the Application (client) ID, Directory (tenant) ID, and client secret value\n\n'
        'Thank you!'
    )
    p = doc.add_paragraph()
    run = p.add_run(email_text)
    run.font.size = Pt(10)
    run.italic = True
    p.paragraph_format.left_indent = Inches(0.5)

    # ── Steps ──
    steps = [
        {
            "title": "Step 1: Open Factory.AI",
            "content": [
                ("text", 'Open your terminal (on Mac: search "Terminal" in Spotlight).'),
                ("text", "Type the following and press Enter:"),
                ("code", "droid"),
                ("text", "You should see the Factory.AI prompt appear. This is where you will give instructions in plain English."),
                ("tip", "Think of Factory.AI as a colleague who writes code. You describe what you want; it builds it."),
            ],
        },
        {
            "title": "Step 2: Tell Factory.AI About Your Dashboard",
            "content": [
                ("text", "Type a message like the one below. Replace the details with your own:"),
                ("code",
                 'Create a Streamlit dashboard using the Excel file at\n'
                 '"/Users/YourName/OneDrive - EY/my_data_file.xlsx"\n'
                 'as the data source. Make it update live when the Excel file changes.\n'
                 'Use EY branding (dark theme, yellow accents). Include KPI cards,\n'
                 'charts, and filters. Anyone with an @ey.com email can view.\n'
                 'Only admins can edit data. Let me add and remove admins\n'
                 'directly inside the app by entering their @ey.com email.\n'
                 'Include Microsoft SSO authentication using Azure AD so users\n'
                 'sign in with their real EY Microsoft account. Fall back to\n'
                 'manual email input if Azure secrets are not yet configured.'),
                ("heading3", "How to find your Excel file path:"),
                ("numbered", [
                    "Open Finder",
                    'Navigate to your Excel file',
                    'Right-click the file > Hold the Option key > Click "Copy as Pathname"',
                    "Paste it into your message above",
                ]),
                ("text", "Factory.AI will read your Excel file, understand the data, and build the full dashboard application."),
                ("text", "Wait for it to finish. It will tell you when it is done."),
            ],
        },
        {
            "title": "Step 3: Test Your Dashboard Locally",
            "content": [
                ("text", "Factory.AI will typically launch the dashboard for you. If not, type:"),
                ("code", "run the streamlit app"),
                ("text", "Your browser will open with your dashboard. Verify:"),
                ("bullets", [
                    "The charts display your data correctly",
                    "The filters work (Region, Date Range, etc.)",
                    "The KPI numbers look reasonable",
                ]),
                ("tip", 'Want changes? Just tell Factory.AI in plain English:\n'
                        '  - "Make the title bigger"\n'
                        '  - "Add a bar chart showing revenue by region"\n'
                        '  - "Change the color scheme"'),
            ],
        },
        {
            "title": "Step 4: Create a GitHub Repository",
            "content": [
                ("text", "Tell Factory.AI:"),
                ("code", 'push this code to GitHub and create a new repo called "My Dashboard Name"'),
                ("text", "Factory.AI will create a new repository on GitHub, upload all the code, and give you the repository URL."),
                ("text", "Save this URL. You will need it in the next step."),
            ],
        },
        {
            "title": "Step 5: Deploy to Streamlit Community Cloud",
            "content": [
                ("text", "This is how you make the dashboard accessible to everyone at EY via a web link."),
                ("heading3", "5a. Sign in to Streamlit Cloud"),
                ("numbered", [
                    "Go to share.streamlit.io",
                    'Click "Continue with GitHub"',
                    "Sign in with your GitHub account",
                    "Authorize Streamlit to access your repositories",
                ]),
                ("heading3", "5b. Deploy Your App"),
                ("numbered", [
                    'Click "New app" (top right)',
                    "Fill in:\n"
                    "    - Repository: Select your dashboard repo\n"
                    "    - Branch: main\n"
                    "    - Main file path: app.py",
                    'Click "Deploy"',
                ]),
                ("heading3", "5c. Wait for Deployment"),
                ("bullets", [
                    "Streamlit will install dependencies and start your app",
                    "This takes 2-3 minutes the first time",
                    "You will see a live URL like: https://my-dashboard-name-xyz.streamlit.app",
                ]),
                ("tip", 'If you see an error: copy the error message, go back to Factory.AI, paste it, and say "fix this error." Then push the fix to GitHub. Streamlit Cloud will automatically redeploy.'),
            ],
        },
        {
            "title": "Step 6: Share With Your Team",
            "content": [
                ("text", "Your dashboard is now live. Share the Streamlit URL with your team. For sensitive data, complete Step 8 (SSO) before sharing broadly."),
                ("heading3", "Two Roles: Viewer vs Admin"),
                ("bullets", [
                    "Viewer (any @ey.com email): View charts, use filters, explore data",
                    "Admin (only people you grant): Edit data, download CSV, save to Excel, manage other admins",
                ]),
                ("text", "You are the first admin. Your email is set as admin by default. Everyone else starts as a viewer."),
                ("heading3", "How to Give Someone Admin Access"),
                ("numbered", [
                    "Open the dashboard and sign in (you must be an admin)",
                    'Scroll to the bottom: "Admin: Manage Admins"',
                    "Enter the person's full @ey.com email in the 'Add admin email' field",
                    'Click "Add Admin" \u2014 they have admin access immediately',
                ]),
                ("heading3", "How to Remove Admin Access"),
                ("numbered", [
                    'Scroll to "Admin: Manage Admins"',
                    "Use the 'Remove admin' dropdown to select the person",
                    'Click "Remove Admin" \u2014 they revert to viewer immediately',
                ]),
                ("tip", "You cannot remove yourself as admin. There must always be at least one admin."),
                ("heading3", "What Each Role Sees"),
                ("bullets", [
                    "Viewers: Charts, KPIs, filters, read-only data table (in expandable section)",
                    "Admins: All of the above PLUS editable data table, save to Excel, CSV download, admin management panel",
                ]),
            ],
        },
        {
            "title": "Step 7: Keep Your Dashboard Updated",
            "content": [
                ("heading3", "To update the data:"),
                ("numbered", [
                    "Open Factory.AI",
                    'Say: "copy the latest version of my Excel file into the data folder and push to GitHub"',
                    "Streamlit Cloud will automatically redeploy with the new data (~1 minute)",
                ]),
                ("heading3", "To change the dashboard:"),
                ("numbered", [
                    "Open Factory.AI",
                    'Describe what you want (e.g., "add a chart showing quarterly trends")',
                    'Say: "push the changes to GitHub"',
                    "Streamlit Cloud picks up the changes automatically",
                ]),
            ],
        },
        {
            "title": "Step 8: Activate Microsoft SSO",
            "content": [
                ("text", "Your dashboard already has SSO built in from Step 2. Now you activate it by pasting the three values your IT team provided (from the request you sent before starting)."),
                ("heading3", "Why This Step Is Essential"),
                ("text", "Without SSO, anyone can type any @ey.com email into the dashboard with no verification. SSO forces users to sign in with their real EY Microsoft account \u2014 the same login they use for Outlook and Teams. This is required before sharing any sensitive data."),
                ("tip", "No code changes needed. You are only pasting configuration values. The SSO code is already in your app."),
                ("heading3", "What You Need"),
                ("text", "If your IT team has already responded to the email you sent at the start of this guide, you should have the three values (client_id, tenant_id, client_secret) ready to go. If not, follow the steps below or share this section with your IT contact."),
                ("heading3", "8a. Register an App in Microsoft Entra ID"),
                ("numbered", [
                    "Go to portal.azure.com and sign in with your EY account",
                    'In the top search bar, type "App registrations" and click on it',
                    'Click "+ New registration"',
                    'Name: "Factory AI Dashboard"',
                    'Supported account types: "Accounts in this organizational directory only"',
                    "Redirect URI: Platform = Web, URI = your Streamlit Cloud URL",
                    'Click "Register"',
                ]),
                ("heading3", "8b. Copy Your Two IDs"),
                ("text", "On the app's Overview page, copy and save:"),
                ("bullets", [
                    "Application (client) ID — save as client_id",
                    "Directory (tenant) ID — save as tenant_id",
                ]),
                ("heading3", "8c. Create a Client Secret"),
                ("numbered", [
                    'In the left sidebar, click "Certificates & secrets"',
                    'Click "+ New client secret"',
                    'Enter a description (e.g., "streamlit-dashboard")',
                    "Choose an expiry period (e.g., 12 months)",
                    'Click "Add"',
                    'Immediately copy the value in the "Value" column — this is your client_secret',
                ]),
                ("tip", "The secret value is only shown once. If you navigate away without copying it, you will need to create a new one."),
                ("heading3", "8d. Add API Permissions"),
                ("numbered", [
                    'In the left sidebar, click "API permissions"',
                    'Click "+ Add a permission"',
                    'Select "Microsoft Graph"',
                    'Select "Delegated permissions"',
                    'Search for "User.Read" and check the box',
                    'Click "Add permissions"',
                    'Click "Grant admin consent for EY"',
                ]),
                ("text", 'If you do not see the "Grant admin consent" button, you need someone with a higher admin role to click it for you.'),
                ("heading3", "8e. Add Your Three Values to Streamlit Cloud"),
                ("text", "This is the final step. No code changes are needed."),
                ("numbered", [
                    "Go to share.streamlit.io",
                    "Find your app, click the three-dot menu > Settings",
                    'Click the "Secrets" tab',
                    "Paste the following (with your actual values):",
                ]),
                ("code",
                 '[azure]\n'
                 'client_id = "paste-your-application-client-id-here"\n'
                 'client_secret = "paste-your-client-secret-value-here"\n'
                 'tenant_id = "paste-your-directory-tenant-id-here"\n'
                 'redirect_uri = "https://your-app-name.streamlit.app"'),
                ("numbered", [
                    'Click "Save"',
                    'Click "Reboot app"',
                ]),
                ("heading3", "8f. Verify It Works"),
                ("numbered", [
                    "Open your dashboard URL in a private/incognito browser window",
                    'You should see a "Sign in with Microsoft" button',
                    "Click it — you will be redirected to Microsoft's login page",
                    "Sign in with your EY account",
                    "You should be returned to the dashboard showing your name and email",
                ]),
                ("heading3", "What Changes for Users"),
                ("bullets", [
                    "Before SSO: Users type their email into a text box — no verification",
                    "After SSO: Users click 'Sign in with Microsoft' — Microsoft verifies their identity",
                    "After SSO: Only suitable for sensitive and confidential data",
                    "After SSO: Microsoft logs all sign-ins for audit purposes",
                ]),
                ("heading3", "Email Template for Your IT Team"),
                ("text",
                 'If you need to request help from IT, send them this:\n\n'
                 '"Hi, I need an Azure AD (Microsoft Entra ID) app registration for an internal '
                 'Streamlit dashboard. Could you help with the following?\n'
                 '1. Register a new app called Factory AI Dashboard\n'
                 '2. Set supported account types to Accounts in this organizational directory only\n'
                 '3. Add a Web redirect URI: [your Streamlit URL]\n'
                 '4. Create a client secret\n'
                 '5. Add the Microsoft Graph User.Read delegated permission and grant admin consent\n'
                 '6. Send me the Application (client) ID, Directory (tenant) ID, and client secret value\n'
                 'Thank you!"'),
            ],
        },
    ]

    for step in steps:
        doc.add_heading(step["title"], level=1)
        for block_type, block in step["content"]:
            if block_type == "text":
                doc.add_paragraph(block)
            elif block_type == "code":
                p = doc.add_paragraph()
                run = p.add_run(block)
                run.font.name = "Consolas"
                run.font.size = Pt(10)
                run.font.color.rgb = EY_DARK
                p.paragraph_format.left_indent = Inches(0.5)
            elif block_type == "tip":
                p = doc.add_paragraph()
                run = p.add_run("Tip: ")
                run.bold = True
                run.font.color.rgb = RGBColor(0x00, 0xA3, 0xE0)
                p.add_run(block)
            elif block_type == "bullets":
                for item in block:
                    doc.add_paragraph(item, style="List Bullet")
            elif block_type == "numbered":
                for i, item in enumerate(block):
                    doc.add_paragraph(item, style="List Number")
            elif block_type == "heading3":
                doc.add_heading(block, level=3)

    # ── Troubleshooting ──
    doc.add_heading("Troubleshooting", level=1)

    trouble = [
        ["Problem", "Solution"],
        ['"ModuleNotFoundError" on Streamlit Cloud', 'Tell Factory.AI: "add [package name] to requirements.txt and push to GitHub"'],
        ['"Excel file not found" on Streamlit Cloud', 'Tell Factory.AI: "copy my Excel file into the data folder in the repo and push"'],
        ["Dashboard looks broken after changes", 'Go back to Factory.AI, paste the error, and say "fix this." Then push to GitHub.'],
        ["Charts show wrong data", "Check your Excel for blank rows, text in numeric columns, or merged cells."],
        ["Can't access Streamlit Cloud", "Sign up at streamlit.io/cloud with the same GitHub account that owns the repo."],
        ['SSO: "AADSTS50011" redirect URI error', "The redirect URI in Azure does not match your Streamlit URL exactly. Fix it in Azure Portal > App registration > Authentication."],
        ['SSO: "Authentication failed"', "The client secret may have expired. Create a new one in Azure Portal and update Streamlit Cloud secrets."],
        ['SSO: "Grant admin consent" button missing', "You need a Global Administrator or Privileged Role Administrator to grant consent. Ask your IT admin."],
    ]
    table = doc.add_table(rows=len(trouble), cols=2)
    table.style = "Light Grid Accent 1"
    for r, row_data in enumerate(trouble):
        for c, val in enumerate(row_data):
            cell = table.rows[r].cells[c]
            cell.text = val
            if r == 0:
                cell.paragraphs[0].runs[0].bold = True

    # ── Quick Reference ──
    doc.add_heading("Quick Reference", level=1)

    qr = [
        ["Task", "What to Tell Factory.AI"],
        ["Build a dashboard", '"Create a Streamlit dashboard using [file path]"'],
        ["Change a chart", '"Change the bar chart to a line chart"'],
        ["Add a filter", '"Add a filter for [column name]"'],
        ["Fix an error", 'Paste the error and say "fix this"'],
        ["Push to GitHub", '"Push the changes to GitHub"'],
    ]
    table = doc.add_table(rows=len(qr), cols=2)
    table.style = "Light Grid Accent 1"
    for r, row_data in enumerate(qr):
        for c, val in enumerate(row_data):
            cell = table.rows[r].cells[c]
            cell.text = val
            if r == 0:
                cell.paragraphs[0].runs[0].bold = True

    # ── Footer ──
    doc.add_paragraph()
    p = doc.add_paragraph()
    p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run = p.add_run("Built with Factory.AI — the AI software engineering platform.")
    run.italic = True
    run.font.color.rgb = EY_GREY

    path = os.path.join(OUT_DIR, "Factory_AI_Dashboard_Guide.docx")
    doc.save(path)
    print(f"Word doc saved: {path}")


# ═══════════════════════════════════════════════════════════════════════════
# POWERPOINT
# ═══════════════════════════════════════════════════════════════════════════

def make_ppt():
    prs = Presentation()
    prs.slide_width = PptInches(13.333)
    prs.slide_height = PptInches(7.5)

    def add_slide(title_text, bullets=None, subtitle_text=None, is_title_slide=False):
        if is_title_slide:
            layout = prs.slide_layouts[6]  # blank
        else:
            layout = prs.slide_layouts[6]  # blank

        slide = prs.slides.add_slide(layout)

        # Dark background
        bg = slide.background
        fill = bg.fill
        fill.solid()
        fill.fore_color.rgb = PPT_DARK

        # Yellow accent bar at top
        from pptx.util import Emu
        shape = slide.shapes.add_shape(
            1, 0, 0, prs.slide_width, PptInches(0.08)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = PPT_YELLOW
        shape.line.fill.background()

        if is_title_slide:
            # Title
            txBox = slide.shapes.add_textbox(
                PptInches(1), PptInches(2.2), PptInches(11.333), PptInches(2)
            )
            tf = txBox.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = title_text
            p.font.size = PptPt(40)
            p.font.color.rgb = PPT_WHITE
            p.font.bold = True
            p.alignment = PP_ALIGN.CENTER

            if subtitle_text:
                p2 = tf.add_paragraph()
                p2.text = subtitle_text
                p2.font.size = PptPt(22)
                p2.font.color.rgb = PPT_YELLOW
                p2.alignment = PP_ALIGN.CENTER
                p2.space_before = PptPt(16)
        else:
            # Title
            txBox = slide.shapes.add_textbox(
                PptInches(0.8), PptInches(0.35), PptInches(11.733), PptInches(0.9)
            )
            tf = txBox.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = title_text
            p.font.size = PptPt(30)
            p.font.color.rgb = PPT_YELLOW
            p.font.bold = True

            # Bullets
            if bullets:
                txBox2 = slide.shapes.add_textbox(
                    PptInches(1), PptInches(1.5), PptInches(11.333), PptInches(5.5)
                )
                tf2 = txBox2.text_frame
                tf2.word_wrap = True
                for i, bullet in enumerate(bullets):
                    if i == 0:
                        p = tf2.paragraphs[0]
                    else:
                        p = tf2.add_paragraph()
                    p.text = bullet
                    p.font.size = PptPt(20)
                    p.font.color.rgb = PPT_WHITE
                    p.space_before = PptPt(10)
                    p.space_after = PptPt(4)

        return slide

    # ── Slide 1: Title ──
    add_slide(
        "Building and Deploying a\nLive Dashboard with Factory.AI",
        subtitle_text="A Step-by-Step Guide for EY Professionals",
        is_title_slide=True,
    )

    # ── Slide 2: What You Will Build ──
    add_slide("What You Will Build", [
        "A live, interactive dashboard from your Excel data",
        "Automatic updates when your spreadsheet changes",
        "Accessible to everyone at EY via a web link",
        "Admin controls for editing and access management",
        "Secured with Microsoft SSO — real EY account sign-in",
        "",
        "Time: ~30 minutes (including SSO)  |  Technical experience: None",
    ])

    # ── Slide 3: What You Need ──
    add_slide("What You Need Before Starting", [
        "1.  Factory.AI CLI (Droid) — pre-installed or visit factory.ai",
        "2.  A GitHub account — github.com (free)",
        "3.  A Streamlit Cloud account — streamlit.io/cloud (free)",
        "4.  Your data in an Excel file (.xlsx)",
        "5.  Azure AD app registration — request from IT (see next slide)",
    ])

    # ── Slide 3b: Send IT Request Now ──
    add_slide("FIRST: Send This Email to Your IT Team", [
        "Send this request NOW — it may take a day or two to process.",
        "The dashboard works without SSO while you wait.",
        "",
        '"Hi, I need an Azure AD app registration for a Streamlit dashboard.',
        '  1. Register app: Factory AI Dashboard',
        '  2. Account type: This organizational directory only',
        '  3. Redirect URI: Web > [I will provide the URL]',
        '  4. Create a client secret',
        '  5. Add User.Read permission + grant admin consent',
        '  6. Send me: client_id, tenant_id, client_secret"',
    ])

    # ── Slide 4: Step 1 ──
    add_slide("Step 1: Open Factory.AI", [
        'Open your terminal (Mac: search "Terminal" in Spotlight)',
        "",
        'Type:  droid',
        "",
        "The Factory.AI prompt will appear.",
        "You give instructions in plain English — it writes the code.",
    ])

    # ── Slide 5: Step 2 ──
    add_slide("Step 2: Describe Your Dashboard", [
        "Tell Factory.AI what you want. Example:",
        "",
        '   "Create a Streamlit dashboard using [my file path].',
        '    Use EY branding. Include KPI cards, charts, and filters.',
        '    Anyone with an @ey.com email can view.',
        '    Only admins can edit. Let me add and remove admins',
        '    directly inside the app by entering their @ey.com email.',
        '    Include Microsoft SSO using Azure AD. Fall back to',
        '    manual email input if Azure secrets are not yet configured."',
        "",
        "Tip: Right-click your Excel file > Option + Copy as Pathname",
    ])

    # ── Slide 6: Step 3 ──
    add_slide("Step 3: Test Locally", [
        "Factory.AI launches the dashboard in your browser.",
        'If not, type:  "run the streamlit app"',
        "",
        "Verify:",
        "   Charts display your data correctly",
        "   Filters work (Region, Date Range, etc.)",
        "   KPI numbers look reasonable",
        "",
        'Want changes? Just say it: "Make the title bigger"',
    ])

    # ── Slide 7: Step 4 ──
    add_slide("Step 4: Push to GitHub", [
        "Tell Factory.AI:",
        "",
        '   "Push this code to GitHub and create a new',
        '    repo called My Dashboard Name"',
        "",
        "It will create the repository and upload everything.",
        "Save the URL it gives you — you need it next.",
    ])

    # ── Slide 8: Step 5 ──
    add_slide("Step 5: Deploy to Streamlit Cloud", [
        "1.  Go to share.streamlit.io and sign in with GitHub",
        '2.  Click "New app"',
        "3.  Select your repository, branch: main, file: app.py",
        '4.  Click "Deploy"',
        "5.  Wait 2-3 minutes for the app to build",
        "",
        "You will get a live URL like:",
        "   https://my-dashboard-xyz.streamlit.app",
    ])

    # ── Slide 9: Step 6 ──
    add_slide("Step 6: Share With Your Team", [
        "Share the Streamlit URL to start testing.",
        "For sensitive data, complete Step 8 (SSO) before sharing broadly.",
        "",
        "Two roles:",
        "   Viewer (any @ey.com email)  →  View charts, filters, read-only data",
        "   Admin (people you grant)      →  Edit data, download, manage access",
        "",
        "You are the first admin. Everyone else starts as a viewer.",
    ])

    # ── Slide 9b: Managing Access ──
    add_slide("Managing Admin vs Viewer Access", [
        "To give someone admin access:",
        "   1.  Sign in as admin > scroll to 'Admin: Manage Admins'",
        "   2.  Enter their @ey.com email > click 'Add Admin'",
        "   3.  They have admin access immediately",
        "",
        "To remove admin access:",
        "   1.  Use the 'Remove admin' dropdown > click 'Remove Admin'",
        "   2.  They revert to viewer access immediately",
        "",
        "Note: You cannot remove yourself. At least one admin is required.",
    ])

    # ── Slide 10: Step 7 ──
    add_slide("Step 7: Keep It Updated", [
        "To update data:",
        '   Tell Factory.AI: "Copy latest Excel into data folder',
        '   and push to GitHub"',
        "",
        "To change the dashboard:",
        '   Describe what you want → "push changes to GitHub"',
        "",
        "Streamlit Cloud redeploys automatically (~1 minute).",
    ])

    # ── Slide 11: SSO - Activate ──
    add_slide("Step 8: Activate Microsoft SSO", [
        "Your app already has SSO built in from Step 2.",
        "Now you activate it with the 3 values from your IT team.",
        "",
        "Without SSO: anyone can type any @ey.com email — no verification.",
        "With SSO: users sign in with their real EY Microsoft account.",
        "",
        "Zero code changes needed — you are only pasting configuration.",
    ])

    # ── Slide 12: SSO - Register the App ──
    add_slide("Step 8a: Register an App in Azure", [
        "You need someone with Azure AD admin access (typically IT).",
        "",
        "1.  Go to portal.azure.com > search 'App registrations'",
        "2.  Click '+ New registration'",
        "3.  Name: 'Factory AI Dashboard'",
        "4.  Account type: 'This organizational directory only' (EY only)",
        "5.  Redirect URI: Web > your Streamlit Cloud URL",
        "6.  Click 'Register'",
    ])

    # ── Slide 13: SSO - Get Credentials ──
    add_slide("Step 8b-c: Get Your Three Values", [
        "From the app Overview page, copy:",
        "   1.  Application (client) ID  →  your client_id",
        "   2.  Directory (tenant) ID    →  your tenant_id",
        "",
        "Then go to Certificates & secrets:",
        "   3.  Click '+ New client secret' > copy the Value  →  your client_secret",
        "",
        "IMPORTANT: The secret is only shown once — copy it immediately.",
    ])

    # ── Slide 14: SSO - Permissions ──
    add_slide("Step 8d: Add API Permissions", [
        "1.  Go to API permissions > '+ Add a permission'",
        "2.  Select 'Microsoft Graph' > 'Delegated permissions'",
        "3.  Search for 'User.Read' and check it",
        "4.  Click 'Add permissions'",
        "5.  Click 'Grant admin consent for EY'",
        "",
        "If you don't see the consent button, ask a Global Admin to do it.",
    ])

    # ── Slide 15: SSO - Paste into Streamlit ──
    add_slide("Step 8e: Paste Values into Streamlit Cloud", [
        "1.  Go to share.streamlit.io > your app > Settings > Secrets",
        "2.  Paste:",
        "",
        '     [azure]',
        '     client_id = "your-client-id"',
        '     client_secret = "your-client-secret"',
        '     tenant_id = "your-tenant-id"',
        '     redirect_uri = "https://your-app.streamlit.app"',
        "",
        "3.  Click Save > Reboot app. That's it — SSO is live.",
    ])

    # ── Slide 16: SSO - Before vs After ──
    add_slide("What Changes After SSO", [
        "BEFORE:                                    AFTER:",
        "  Users type email in a text box              Users click 'Sign in with Microsoft'",
        "  No verification at all                           Microsoft verifies their identity",
        "  Not safe for sensitive data                  Enterprise-grade security",
        "  No audit trail                                      Microsoft logs all sign-ins",
        "",
        "The app switches automatically — no code changes needed.",
    ])

    # ── Slide 17: Reminder ──
    add_slide("Reminder: Follow Up with IT", [
        "If your IT team has not yet responded to the email",
        "you sent at the beginning of this guide:",
        "",
        "  1.  Follow up and provide your Streamlit Cloud URL",
        "       (from Step 5) as the redirect URI",
        "  2.  Share the AZURE_SSO_SETUP.md file from your",
        "       GitHub repository for full technical details",
        "",
        "Once you have the 3 values, paste them into Streamlit",
        "Cloud Secrets (Step 8e) and SSO goes live instantly.",
    ])

    # ── Slide 18: Troubleshooting ──
    add_slide("Troubleshooting", [
        'ModuleNotFoundError  →  "Add [package] to requirements.txt and push"',
        'Excel file not found  →  "Copy Excel into data folder and push"',
        'Dashboard broken  →  Paste error, say "fix this", then push',
        "Wrong data  →  Check Excel for blank rows or text in number columns",
        'SSO redirect error  →  Verify redirect URI matches exactly in Azure',
        'SSO auth failed  →  Client secret may have expired; create a new one',
    ])

    # ── Slide 19: Quick Reference ──
    add_slide("Quick Reference: What to Tell Factory.AI", [
        'Build a dashboard     →  "Create a Streamlit dashboard using [file]"',
        'Change a chart          →  "Change the bar chart to a line chart"',
        'Add a filter                 →  "Add a filter for [column name]"',
        'Fix an error                →  Paste the error, say "fix this"',
        'Push to GitHub          →  "Push the changes to GitHub"',
    ])

    # ── Slide 20: Closing ──
    add_slide(
        "Start Building Today",
        subtitle_text="factory.ai",
        is_title_slide=True,
    )

    path = os.path.join(OUT_DIR, "Factory_AI_Dashboard_Guide.pptx")
    prs.save(path)
    print(f"PowerPoint saved: {path}")


if __name__ == "__main__":
    make_word()
    make_ppt()
